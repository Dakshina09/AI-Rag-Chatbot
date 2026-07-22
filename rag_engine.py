"""
rag_engine.py
Core RAG (Retrieval-Augmented Generation) logic for the College Notes Chatbot.

Handles:
  - Loading text out of PDF / DOCX / PPTX / TXT files
  - Chunking text into overlapping windows
  - Embedding chunks locally with sentence-transformers (no API key needed)
  - Storing/searching vectors with FAISS, persisted to disk per "subject"
  - Calling an LLM (OpenAI or Anthropic) to generate a grounded answer,
    with a safe extractive fallback if no API key is configured
"""

import pickle
from pathlib import Path

import numpy as np
import requests
import faiss
from sentence_transformers import SentenceTransformer
import pypdf
import docx
from pptx import Presentation

DATA_DIR = Path(__file__).parent / "rag_data"
DATA_DIR.mkdir(exist_ok=True)


# --------------------------------------------------------------------------
# 1. Document loading
# --------------------------------------------------------------------------
class DocumentLoader:
    """Extracts raw text from common note formats."""

    @staticmethod
    def load(file_path):
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return DocumentLoader._load_pdf(file_path)
        elif ext == ".docx":
            return DocumentLoader._load_docx(file_path)
        elif ext == ".pptx":
            return DocumentLoader._load_pptx(file_path)
        elif ext == ".txt":
            return DocumentLoader._load_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    @staticmethod
    def _load_pdf(file_path):
        text = ""
        reader = pypdf.PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            text += f"\n[Page {i + 1}]\n{page_text}"
        return text

    @staticmethod
    def _load_docx(file_path):
        d = docx.Document(file_path)
        return "\n".join(p.text for p in d.paragraphs)

    @staticmethod
    def _load_pptx(file_path):
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            text += f"\n[Slide {i + 1}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    def _load_txt(file_path):
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


# --------------------------------------------------------------------------
# 2. Chunking
# --------------------------------------------------------------------------
def chunk_text(text, chunk_size=800, overlap=150):
    """Simple sliding-window character chunker with overlap so context
    isn't lost at chunk boundaries."""
    chunks = []
    text = text.strip()
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks


# --------------------------------------------------------------------------
# 3. Vector store (FAISS) - one persistent index per "subject" notebook
# --------------------------------------------------------------------------
class VectorStore:
    def __init__(self, subject, embed_model="all-MiniLM-L6-v2"):
        self.subject = subject
        self.dir = DATA_DIR / subject
        self.dir.mkdir(parents=True, exist_ok=True)
        self.index_path = self.dir / "index.faiss"
        self.meta_path = self.dir / "meta.pkl"
        self.embedder = SentenceTransformer(embed_model)
        self.dim = self.embedder.get_sentence_embedding_dimension()
        self.index = None
        self.metadata = []  # parallel list: {text, source, chunk_id}
        self._load()

    def _load(self):
        if self.index_path.exists() and self.meta_path.exists():
            self.index = faiss.read_index(str(self.index_path))
            with open(self.meta_path, "rb") as f:
                self.metadata = pickle.load(f)
        else:
            self.index = faiss.IndexFlatIP(self.dim)  # cosine sim via normalized vectors
            self.metadata = []

    def _save(self):
        faiss.write_index(self.index, str(self.index_path))
        with open(self.meta_path, "wb") as f:
            pickle.dump(self.metadata, f)

    def add_document(self, text, source_name):
        chunks = chunk_text(text)
        if not chunks:
            return 0
        embeddings = self.embedder.encode(chunks, normalize_embeddings=True)
        embeddings = np.array(embeddings, dtype="float32")
        self.index.add(embeddings)
        for i, c in enumerate(chunks):
            self.metadata.append({"text": c, "source": source_name, "chunk_id": i})
        self._save()
        return len(chunks)

    def search(self, query, top_k=4):
        if self.index.ntotal == 0:
            return []
        q_emb = self.embedder.encode([query], normalize_embeddings=True)
        q_emb = np.array(q_emb, dtype="float32")
        scores, idxs = self.index.search(q_emb, min(top_k, self.index.ntotal))
        results = []
        for score, idx in zip(scores[0], idxs[0]):
            if idx == -1:
                continue
            meta = self.metadata[idx]
            results.append({"text": meta["text"], "source": meta["source"], "score": float(score)})
        return results

    def list_sources(self):
        return sorted(set(m["source"] for m in self.metadata))

    def remove_source(self, source_name):
        """Rebuild the index without the given source file."""
        keep = [m for m in self.metadata if m["source"] != source_name]
        self.index = faiss.IndexFlatIP(self.dim)
        self.metadata = []
        if keep:
            texts = [m["text"] for m in keep]
            embeddings = np.array(self.embedder.encode(texts, normalize_embeddings=True), dtype="float32")
            self.index.add(embeddings)
            self.metadata = keep
        self._save()

    def clear(self):
        self.index = faiss.IndexFlatIP(self.dim)
        self.metadata = []
        self._save()

    @staticmethod
    def list_subjects():
        if not DATA_DIR.exists():
            return []
        return sorted(p.name for p in DATA_DIR.iterdir() if p.is_dir())


# --------------------------------------------------------------------------
# 4. Prompt construction + LLM calls
# --------------------------------------------------------------------------
def build_prompt(query, retrieved_chunks, history=None):
    context = "\n\n---\n\n".join(
        f"[Source: {c['source']}]\n{c['text']}" for c in retrieved_chunks
    ) or "(no relevant notes found)"

    history_block = ""
    if history:
        history_block = "\n".join(f"{role.upper()}: {msg}" for role, msg in history[-6:])
        history_block = f"\nRECENT CONVERSATION:\n{history_block}\n"

    return f"""You are a helpful study assistant answering questions using the student's own college notes.
Use ONLY the context below to answer. If the answer isn't fully covered, say so honestly and note what's missing.
Keep answers clear and exam-friendly (use short paragraphs or bullet points where helpful).
End your answer with a line listing the source file(s) you used.
{history_block}
CONTEXT FROM NOTES:
{context}

QUESTION: {query}

ANSWER:"""


def call_gemini(prompt, api_key, model="gemini-2.5-flash"):
    """Google's Gemini API free tier -- no credit card needed, and Flash models
    handle very long note documents well thanks to a large context window."""
    resp = requests.post(
        f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={api_key}",
        headers={"Content-Type": "application/json"},
        json={"contents": [{"parts": [{"text": prompt}]}]},
        timeout=60,
    )
    if not resp.ok:
        raise RuntimeError(f"Gemini API error {resp.status_code}: {resp.text}")
    data = resp.json()
    return data["candidates"][0]["content"]["parts"][0]["text"]


def call_groq(prompt, api_key, model="llama-3.3-70b-versatile"):
    """Groq has a free developer tier (no credit card needed), and hosts several
    open models -- great for a no-cost setup. OpenAI-compatible endpoint."""
    resp = requests.post(
        "https://api.groq.com/openai/v1/chat/completions",
        headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
        json={
            "model": model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "max_tokens": 900,
        },
        timeout=60,
    )
    if not resp.ok:
        raise RuntimeError(f"Groq API error {resp.status_code}: {resp.text}")
    return resp.json()["choices"][0]["message"]["content"]


def extractive_fallback(retrieved_chunks):
    """Used when no LLM API key is configured -- still useful, just not
    a synthesized answer, so the app works out of the box."""
    if not retrieved_chunks:
        return ("I couldn't find anything relevant in your notes for this question. "
                 "Try uploading notes that cover this topic.")
    out = ("_No AI model is connected right now, so here are the most relevant "
           "passages found directly in your notes:_\n\n")
    for c in retrieved_chunks:
        out += f"**From {c['source']} (relevance {c['score']:.2f}):**\n{c['text']}\n\n"
    return out


def generate_answer(query, retrieved_chunks, provider="none", api_key=None, model=None, history=None):
    prompt = build_prompt(query, retrieved_chunks, history=history)
    if provider == "groq" and api_key:
        return call_groq(prompt, api_key, model or "llama-3.3-70b-versatile")
    elif provider == "gemini" and api_key:
        return call_gemini(prompt, api_key, model or "gemini-2.5-flash")
    else:
        return extractive_fallback(retrieved_chunks)
